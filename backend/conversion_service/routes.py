"""
Conversion Service: PDF <-> Word / Excel / PPT / Image
Uses LibreOffice for Office formats, PyMuPDF for image export, Pandoc for text.
"""
import io
import os
import re
import subprocess
import uuid
from uuid import UUID

import fitz
import httpx
from fastapi import APIRouter, Depends, HTTPException, Header, UploadFile, File, Form
from pydantic import BaseModel
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from shared.s3 import download_file, upload_file, generate_presigned_url
from shared.database import get_db
from shared.malware import scan_bytes, MalwareDetected

router = APIRouter()

SUPPORTED_FORMATS = {"pdf", "docx", "xlsx", "pptx", "png", "jpg", "txt", "html"}
LIBREOFFICE_BIN   = os.getenv("LIBREOFFICE_BIN", "soffice")


async def current_user_id(x_user_id: str | None = Header(None)) -> UUID:
    if not x_user_id:
        raise HTTPException(status_code=401, detail="Not authenticated")
    return UUID(x_user_id)


class ConvertRequest(BaseModel):
    document_id:   str
    source_format: str = "pdf"
    target_format: str
    s3_key:        str | None = None   # ignored; resolved server-side from document_id


class ConvertResponse(BaseModel):
    output_key:   str
    download_url: str


@router.get("/formats")
async def supported_formats():
    return {"formats": sorted(SUPPORTED_FORMATS)}


@router.post("/convert", response_model=ConvertResponse)
async def convert_document(
    body:    ConvertRequest,
    user_id: UUID = Depends(current_user_id),
    db:      AsyncSession = Depends(get_db),
):
    if body.target_format not in SUPPORTED_FORMATS:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {body.target_format}")

    # Resolve the storage key from the DB and enforce ownership — never trust a
    # client-supplied s3_key (that would allow converting other users' documents).
    row = (await db.execute(
        text("SELECT s3_key FROM documents WHERE id = CAST(:id AS uuid) AND owner_id = CAST(:uid AS uuid)"),
        {"id": body.document_id, "uid": str(user_id)},
    )).first()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    source_key = row[0]

    job_id  = str(uuid.uuid4())
    tmp_in  = f"/tmp/conv_{job_id}_in.{body.source_format}"
    download_file(source_key, tmp_in)

    if body.target_format in ("png", "jpg"):
        output_key, tmp_out = _pdf_to_image(tmp_in, job_id, body.target_format, user_id)
    elif body.target_format == "txt":
        output_key, tmp_out = _pdf_to_text(tmp_in, job_id, user_id)
    elif body.target_format == "html":
        output_key, tmp_out = _pdf_to_html(tmp_in, job_id, user_id)
    elif body.source_format == "pdf" and body.target_format == "docx":
        output_key, tmp_out = _pdf_to_docx(tmp_in, job_id, user_id)
    elif body.source_format == "pdf" and body.target_format == "xlsx":
        output_key, tmp_out = _pdf_to_xlsx(tmp_in, job_id, user_id)
    elif body.source_format == "pdf" and body.target_format == "pptx":
        output_key, tmp_out = _pdf_to_pptx(tmp_in, job_id, user_id)
    elif body.target_format == "pdf":
        output_key, tmp_out = _office_to_pdf(tmp_in, job_id, body.source_format, user_id)
    else:
        raise HTTPException(status_code=400, detail="Conversion path not supported")

    with open(tmp_out, "rb") as f:
        mime = _mime(body.target_format)
        upload_file(output_key, f, content_type=mime)

    url = generate_presigned_url(output_key, expiry=3600)
    return ConvertResponse(output_key=output_key, download_url=url)


@router.post("/file", response_model=ConvertResponse)
async def convert_uploaded_file(
    file: UploadFile = File(...),
    target_format: str = Form("pdf"),
    user_id: UUID = Depends(current_user_id),
):
    """Convert an arbitrary uploaded file (Office/image → PDF, or PDF → other).
    Used by the Word/Excel/PPT/Image → PDF tools — bypasses the PDF-only document store."""
    if target_format not in SUPPORTED_FORMATS:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {target_format}")
    ext = (file.filename or "").rsplit(".", 1)[-1].lower() if "." in (file.filename or "") else ""
    job_id = str(uuid.uuid4())
    tmp_in = f"/tmp/conv_{job_id}_in.{ext or 'bin'}"
    payload = await file.read()
    try:
        scan_bytes(payload)
    except MalwareDetected as e:
        raise HTTPException(status_code=400, detail=f"File failed the malware scan ({e.signature})")
    with open(tmp_in, "wb") as f:
        f.write(payload)

    office = {"doc", "docx", "odt", "rtf", "xls", "xlsx", "ods", "csv", "ppt", "pptx", "odp",
              "html", "htm", "txt"}   # LibreOffice renders HTML/text to PDF too
    images = {"png", "jpg", "jpeg", "webp", "bmp", "tiff", "tif", "gif"}

    if target_format == "pdf" and ext in office:
        output_key, tmp_out = _office_to_pdf(tmp_in, job_id, ext, user_id)
    elif target_format == "pdf" and ext in images:
        output_key, tmp_out = _image_to_pdf(tmp_in, job_id, user_id)
    elif ext == "pdf" and target_format in ("png", "jpg"):
        output_key, tmp_out = _pdf_to_image(tmp_in, job_id, target_format, user_id)
    elif ext == "pdf" and target_format == "txt":
        output_key, tmp_out = _pdf_to_text(tmp_in, job_id, user_id)
    else:
        raise HTTPException(status_code=400, detail=f"Cannot convert .{ext or '?'} → {target_format}")

    with open(tmp_out, "rb") as f:
        upload_file(output_key, f, content_type=_mime(target_format))
    return ConvertResponse(output_key=output_key, download_url=generate_presigned_url(output_key, expiry=3600))


@router.post("/protect", response_model=ConvertResponse)
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...),
                      user_id: UUID = Depends(current_user_id)):
    """Encrypt a PDF with a password (AES-256)."""
    if not password.strip():
        raise HTTPException(status_code=400, detail="Password is required")
    job = str(uuid.uuid4())
    tin = f"/tmp/prot_{job}.pdf"
    with open(tin, "wb") as f:
        f.write(await file.read())
    pdf = fitz.open(tin)
    out = f"/tmp/prot_{job}_out.pdf"
    perm = int(fitz.PDF_PERM_ACCESSIBILITY | fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY | fitz.PDF_PERM_ANNOTATE)
    pdf.save(out, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, user_pw=password, permissions=perm)
    pdf.close()
    key = f"conversions/{user_id}/{job}_protected.pdf"
    with open(out, "rb") as f:
        upload_file(key, f, content_type="application/pdf")
    return ConvertResponse(output_key=key, download_url=generate_presigned_url(key, expiry=3600))


@router.post("/unlock", response_model=ConvertResponse)
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(""),
                     user_id: UUID = Depends(current_user_id)):
    """Remove password/encryption from a PDF (needs the current password if set)."""
    job = str(uuid.uuid4())
    tin = f"/tmp/unlock_{job}.pdf"
    with open(tin, "wb") as f:
        f.write(await file.read())
    pdf = fitz.open(tin)
    if pdf.needs_pass:
        if not pdf.authenticate(password):
            pdf.close()
            raise HTTPException(status_code=400, detail="Incorrect password")
    out = f"/tmp/unlock_{job}_out.pdf"
    pdf.save(out, encryption=fitz.PDF_ENCRYPT_NONE)
    pdf.close()
    key = f"conversions/{user_id}/{job}_unlocked.pdf"
    with open(out, "rb") as f:
        upload_file(key, f, content_type="application/pdf")
    return ConvertResponse(output_key=key, download_url=generate_presigned_url(key, expiry=3600))


# ── Translate (English ⇄ Hindi, Telugu, …) ────────────────────────────────────
# Primary: self-hosted LibreTranslate (compose service, no API keys, offline).
# Fallback: the free Google Translate web endpoint (blocked on some networks).

LIBRETRANSLATE_URL = os.getenv("LIBRETRANSLATE_URL", "")

# Offline via LibreTranslate models: en, hi, bn, ur, fr, de, es.
# te/ta have no offline model (not in the Argos catalogue) → online fallback only.
TRANSLATE_LANGS = {
    "en": "English", "hi": "Hindi", "bn": "Bengali", "ur": "Urdu",
    "te": "Telugu", "ta": "Tamil",
    "fr": "French", "de": "German", "es": "Spanish",
}
_TRANSLATE_CHAR_CAP = 30000   # keep interactive latency sane on huge documents


async def _translate_chunk(client: httpx.AsyncClient, text: str, source: str, target: str) -> str:
    if LIBRETRANSLATE_URL:
        try:
            r = await client.post(f"{LIBRETRANSLATE_URL}/translate",
                                  json={"q": text, "source": source, "target": target, "format": "text"})
            r.raise_for_status()
            return r.json()["translatedText"]
        except Exception:
            pass   # fall through to the public endpoint
    r = await client.get(
        "https://translate.googleapis.com/translate_a/single",
        params={"client": "gtx", "sl": source, "tl": target, "dt": "t", "q": text},
    )
    r.raise_for_status()
    if "text/html" in r.headers.get("content-type", ""):
        raise HTTPException(status_code=502, detail="Translation service unavailable")
    return "".join(seg[0] for seg in r.json()[0] if seg and seg[0])


async def _translate_texts(client: httpx.AsyncClient, texts: list[str], source: str, target: str) -> list[str]:
    """Translate many small strings. LibreTranslate accepts a list in one call; batch by
    size so requests stay snappy. Falls back to per-string translation."""
    out: list[str] = []
    batch: list[str] = []
    size = 0

    async def flush():
        nonlocal batch, size
        if not batch:
            return
        if LIBRETRANSLATE_URL:
            try:
                r = await client.post(f"{LIBRETRANSLATE_URL}/translate",
                                      json={"q": batch, "source": source, "target": target, "format": "text"})
                r.raise_for_status()
                tt = r.json()["translatedText"]
                if isinstance(tt, list) and len(tt) == len(batch):
                    out.extend(tt); batch, size = [], 0
                    return
            except Exception:
                pass
        for t in batch:
            out.append(await _translate_chunk(client, t, source, target))
        batch, size = [], 0

    for t in texts:
        if size + len(t) > 6000 or len(batch) >= 80:
            await flush()
        batch.append(t); size += len(t)
    await flush()
    # zero-width joiners/marks from MT render as tofu boxes in many Indic fonts
    return [re.sub("[\\u200b-\\u200f\\u2060\\ufeff]", "", t) for t in out]


# Unicode fonts for re-inserting translated text (first existing candidate wins).
# IMPORTANT: translated labels routinely mix scripts ("Emp Code: AHMS0560" → half
# Latin, half Indic), so a font must cover BOTH — Noto per-script fonts have no
# Latin letters (they render as tofu/NULs); Lohit Devanagari/Telugu and FreeSerif do.
_FONT_CANDIDATES = {
    "hi": ["/usr/share/fonts/truetype/lohit-devanagari/Lohit-Devanagari.ttf",
           "/usr/share/fonts/truetype/freefont/FreeSerif.ttf"],
    "te": ["/usr/share/fonts/truetype/lohit-telugu/Lohit-Telugu.ttf",
           "/usr/share/fonts/truetype/freefont/FreeSerif.ttf"],
    "ta": ["/usr/share/fonts/truetype/freefont/FreeSerif.ttf",
           "/usr/share/fonts/truetype/lohit-tamil/Lohit-Tamil.ttf"],
    "bn": ["/usr/share/fonts/truetype/freefont/FreeSerif.ttf",
           "/usr/share/fonts/truetype/lohit-bengali/Lohit-Bengali.ttf"],
    "ur": ["/usr/share/fonts/truetype/freefont/FreeSerif.ttf",
           "/usr/share/fonts/truetype/noto/NotoNaskhArabic-Regular.ttf"],
}
_DEFAULT_FONT = "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf"


def _font_for(lang: str) -> str:
    for path in _FONT_CANDIDATES.get(lang, []):
        if os.path.exists(path):
            return path
    return _DEFAULT_FONT


def _cluster_unit(spans) -> tuple["fitz.Rect", str, float]:
    rect = fitz.Rect(spans[0]["bbox"])
    for s in spans[1:]:
        rect |= fitz.Rect(s["bbox"])
    text = "".join(s["text"] for s in spans).strip()
    size = max(s["size"] for s in spans)
    return (rect, text, size)


_CELL_GAP_PT = 12.0   # horizontal gap that separates table cells / form columns


def _page_text_blocks(page) -> list[tuple["fitz.Rect", str, float]]:
    """Fine-grained text units: one per line, split at wide horizontal gaps so each
    table cell / form field keeps its own box. Whole-block extraction merged entire
    payslip rows into one clump, destroying the column layout."""
    units = []
    for b in page.get_text("dict")["blocks"]:
        if b.get("type") != 0:
            continue
        for line in b.get("lines", []):
            spans = [s for s in line.get("spans", []) if s["text"].strip()]
            if not spans:
                continue
            spans.sort(key=lambda s: s["bbox"][0])
            cluster = [spans[0]]
            for s in spans[1:]:
                if s["bbox"][0] - cluster[-1]["bbox"][2] > _CELL_GAP_PT:
                    units.append(_cluster_unit(cluster))
                    cluster = [s]
                else:
                    cluster.append(s)
            units.append(_cluster_unit(cluster))
    return [u for u in units if u[1]]


def _insert_fitted(page, rect: "fitz.Rect", text: str, fontfile: str, start_size: float) -> None:
    """Insert text into rect, shrinking until it fits (translations are often longer)."""
    # slight padding: Indic glyph stacks are taller than the Latin line they replace
    rect = fitz.Rect(rect.x0, rect.y0 - 0.5, rect.x1 + 2.0, rect.y1 + 2.0)
    size = min(start_size, 24.0)
    while size >= 4.5:
        rc = page.insert_textbox(rect, text, fontsize=size, fontname="xlat",
                                 fontfile=fontfile, lineheight=1.08)
        if rc >= 0:
            return
        size -= 0.5
    # last resort: tiny font in a slightly grown box so nothing silently disappears
    grown = fitz.Rect(rect.x0, rect.y0, rect.x1 + 20, rect.y1 + 14)
    page.insert_textbox(grown, text, fontsize=4.5, fontname="xlat", fontfile=fontfile, lineheight=1.05)


@router.post("/translate-file")
async def translate_file(
    file: UploadFile = File(...),
    target_lang: str = Form(...),
    source_lang: str = Form("auto"),
    user_id: UUID = Depends(current_user_id),
):
    """Translate a PDF's text into another language. Returns the translated text plus
    a downloadable PDF (rendered with Unicode fonts) and a .txt file."""
    if target_lang not in TRANSLATE_LANGS:
        raise HTTPException(status_code=400, detail=f"target_lang must be one of {sorted(TRANSLATE_LANGS)}")
    if source_lang != "auto" and source_lang not in TRANSLATE_LANGS:
        raise HTTPException(status_code=400, detail="unsupported source_lang")

    data = await file.read()
    if not data.startswith(b"%PDF-"):
        raise HTTPException(status_code=400, detail="Upload a PDF file")
    pdf = fitz.open(stream=data, filetype="pdf")

    # Collect text blocks per page so the translated PDF keeps the ORIGINAL layout —
    # each block is translated and written back into its own box; tables, lines,
    # logos and images stay exactly where they were.
    page_blocks = [_page_text_blocks(page) for page in pdf]
    if not any(page_blocks):
        pdf.close()
        raise HTTPException(status_code=400, detail="No extractable text — run OCR first for scanned documents")

    # translate only blocks that contain letters (amounts/IDs pass through unchanged)
    todo: list[tuple[int, int]] = []
    texts: list[str] = []
    total = 0
    truncated = False
    for pno, blocks in enumerate(page_blocks):
        for bno, (_, text, _) in enumerate(blocks):
            if not any(c.isalpha() for c in text):
                continue
            if total + len(text) > _TRANSLATE_CHAR_CAP:
                truncated = True
                break
            todo.append((pno, bno))
            texts.append(text)
            total += len(text)
        if truncated:
            break

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            translations = await _translate_texts(client, texts, source_lang, target_lang)
    except HTTPException:
        raise
    except Exception:
        pdf.close()
        raise HTTPException(status_code=502, detail="Translation service unavailable — try again shortly")

    translated_map = dict(zip(todo, translations))
    fontfile = _font_for(target_lang)

    for pno, blocks in enumerate(page_blocks):
        page = pdf[pno]
        if not blocks:
            continue
        # white-out original text, keep images and line art
        for rect, _, _ in blocks:
            page.add_redact_annot(rect)
        try:
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE,
                                  graphics=fitz.PDF_REDACT_LINE_ART_NONE)
        except TypeError:   # older PyMuPDF without the graphics kwarg
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
        for bno, (rect, original, size) in enumerate(blocks):
            new_text = translated_map.get((pno, bno), original)
            _insert_fitted(page, rect, new_text, fontfile, size)

    job = str(uuid.uuid4())
    out_pdf = pdf.tobytes(garbage=3, deflate=True)
    pdf.close()

    pdf_key = f"conversions/{user_id}/{job}_translated_{target_lang}.pdf"
    upload_file(pdf_key, io.BytesIO(out_pdf), content_type="application/pdf")

    translated_joined = "\n".join(translations)
    txt_key = f"conversions/{user_id}/{job}_translated_{target_lang}.txt"
    upload_file(txt_key, io.BytesIO(translated_joined.encode("utf-8")), content_type="text/plain; charset=utf-8")

    return {
        "target_language": TRANSLATE_LANGS[target_lang],
        "translated_text": translated_joined[:8000],
        "truncated": truncated,
        "download_url": generate_presigned_url(pdf_key, expiry=3600),
        "txt_url": generate_presigned_url(txt_key, expiry=3600),
    }


@router.get("/translate-languages")
async def translate_languages():
    return {"languages": TRANSLATE_LANGS}


@router.post("/scan", response_model=ConvertResponse)
async def scan_to_pdf(files: list[UploadFile] = File(...), user_id: UUID = Depends(current_user_id)):
    """Combine photos/scanned images into a single PDF (one image per page)."""
    if not files:
        raise HTTPException(status_code=400, detail="Add at least one image")
    job = str(uuid.uuid4())
    out_pdf = fitz.open()
    for idx, f in enumerate(files):
        data = await f.read()
        ext = (f.filename or "").rsplit(".", 1)[-1].lower() or "png"
        tmp_img = f"/tmp/scan_{job}_{idx}.{ext}"
        with open(tmp_img, "wb") as fh:
            fh.write(data)
        img = fitz.open(tmp_img)
        img_pdf = fitz.open("pdf", img.convert_to_pdf())
        out_pdf.insert_pdf(img_pdf)
        img.close(); img_pdf.close()
    out = f"/tmp/scan_{job}_out.pdf"
    out_pdf.save(out)
    out_pdf.close()
    key = f"conversions/{user_id}/{job}_scan.pdf"
    with open(out, "rb") as f:
        upload_file(key, f, content_type="application/pdf")
    return ConvertResponse(output_key=key, download_url=generate_presigned_url(key, expiry=3600))


# ── Conversion Engines ────────────────────────────────────────────────────────

def _image_to_pdf(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    img = fitz.open(src)
    pdf_bytes = img.convert_to_pdf()
    img.close()
    out = f"/tmp/conv_{job_id}.pdf"
    with open(out, "wb") as f:
        f.write(pdf_bytes)
    return f"conversions/{user_id}/{job_id}.pdf", out


def _pdf_to_image(src: str, job_id: str, fmt: str, user_id: UUID) -> tuple[str, str]:
    pdf     = fitz.open(src)
    images  = []
    for i, page in enumerate(pdf):
        pix      = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = f"/tmp/conv_{job_id}_p{i + 1}.{fmt}"
        pix.save(img_path)
        images.append(img_path)
    pdf.close()

    if len(images) == 1:
        out = images[0]
    else:
        from PIL import Image
        imgs = [Image.open(p) for p in images]
        out  = f"/tmp/conv_{job_id}_all.{fmt}"
        imgs[0].save(out, save_all=True, append_images=imgs[1:])

    key = f"conversions/{user_id}/{job_id}.{fmt}"
    return key, out


def _pdf_to_text(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    pdf  = fitz.open(src)
    text = "\n\n".join(page.get_text() for page in pdf)
    pdf.close()

    out = f"/tmp/conv_{job_id}.txt"
    with open(out, "w", encoding="utf-8") as f:
        f.write(text)

    key = f"conversions/{user_id}/{job_id}.txt"
    return key, out


def _pdf_to_html(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    pdf   = fitz.open(src)
    pages = [page.get_text("html") for page in pdf]
    pdf.close()

    html = f"<html><body>{''.join(pages)}</body></html>"
    out  = f"/tmp/conv_{job_id}.html"
    with open(out, "w", encoding="utf-8") as f:
        f.write(html)

    key = f"conversions/{user_id}/{job_id}.html"
    return key, out


def _pdf_to_docx(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    # pdf2docx reconstructs paragraphs/tables/layout — far more reliable than LibreOffice
    # for the PDF→Word direction (which LibreOffice silently produces no output for).
    from pdf2docx import Converter
    out = f"/tmp/conv_{job_id}.docx"
    cv = Converter(src)
    try:
        cv.convert(out)
    finally:
        cv.close()
    return f"conversions/{user_id}/{job_id}.docx", out


def _coerce(v):
    """Store numbers as real numbers (so totals/sorting work in Excel), text as text.
    Handles thousands separators and currency-ish values like '1,800.00'."""
    if v is None:
        return None
    s = str(v).strip()
    num = s.replace(",", "").replace("₹", "").replace("$", "").strip()
    try:
        if num and num.lstrip("-").replace(".", "", 1).isdigit():
            f = float(num)
            return int(f) if f.is_integer() else f
    except ValueError:
        pass
    return s


def _page_grid(page, y_tol: float = 4.0, cell_gap: float = 11.0, col_gap: float = 18.0):
    """Reconstruct a page's visual grid from word positions so Excel mirrors the PDF.
      • words on the same baseline (within y_tol) form a row,
      • a horizontal gap > cell_gap starts a new cell,
      • cell start-X values are clustered (col_gap) into shared columns so rows line up.
    Returns (rows, ncols) where rows is a list of {col_index: text}."""
    words = page.get_text("words")          # (x0, y0, x1, y1, word, block, line, wno)
    if not words:
        return [], 0
    words.sort(key=lambda w: (round(w[1] / y_tol), w[0]))

    # group into rows by y
    rows_words, cur, cur_y = [], [], None
    for w in words:
        if cur_y is None or abs(w[1] - cur_y) <= y_tol:
            cur.append(w); cur_y = w[1] if cur_y is None else cur_y
        else:
            rows_words.append(cur); cur, cur_y = [w], w[1]
    if cur:
        rows_words.append(cur)

    # split each row into cells by x-gaps → list of (start_x, text)
    raw_rows = []
    for r in rows_words:
        r.sort(key=lambda w: w[0])
        cells, bucket, last_x1 = [], [], None
        for w in r:
            if last_x1 is not None and (w[0] - last_x1) > cell_gap:
                cells.append((bucket[0][0], " ".join(x[4] for x in bucket))); bucket = []
            bucket.append(w); last_x1 = w[2]
        if bucket:
            cells.append((bucket[0][0], " ".join(x[4] for x in bucket)))
        raw_rows.append(cells)

    # derive shared column start positions from all cell start-X values
    starts = sorted(c[0] for row in raw_rows for c in row)
    col_starts = []
    for x in starts:
        if not col_starts or x - col_starts[-1] > col_gap:
            col_starts.append(x)

    def col_of(x: float) -> int:
        idx = 0
        for i, cx in enumerate(col_starts):
            if x >= cx - 1:
                idx = i
        return idx

    grid = []
    for cells in raw_rows:
        row_map: dict[int, str] = {}
        for sx, txtv in cells:
            ci = col_of(sx)
            row_map[ci] = (row_map[ci] + " " + txtv).strip() if ci in row_map else txtv
        grid.append(row_map)
    return grid, len(col_starts)


def _pdf_to_xlsx(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter

    pdf = fitz.open(src)
    wb = Workbook()
    wb.remove(wb.active)
    for pi, page in enumerate(pdf):
        ws = wb.create_sheet(title=f"Page {pi + 1}"[:31])
        grid, ncols = _page_grid(page)
        widths: dict[int, int] = {}
        for ri, row_map in enumerate(grid, start=1):
            for ci, val in row_map.items():
                ws.cell(row=ri, column=ci + 1, value=_coerce(val))
                widths[ci + 1] = max(widths.get(ci + 1, 10), min(70, len(str(val)) + 2))
        for ci, w in widths.items():
            ws.column_dimensions[get_column_letter(ci)].width = w
        if not grid:   # scanned/image page with no extractable words
            ws.cell(row=1, column=1, value="(no extractable text on this page)")
    pdf.close()
    out = f"/tmp/conv_{job_id}.xlsx"
    wb.save(out)
    return f"conversions/{user_id}/{job_id}.xlsx", out


def _pdf_to_pptx(src: str, job_id: str, user_id: UUID) -> tuple[str, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pdf = fitz.open(src)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for page in pdf:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.3), prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        lines = (page.get_text().strip() or "(no extractable text)").splitlines()
        tf.text = lines[0] if lines else ""
        for line in lines[1:]:
            tf.add_paragraph().text = line
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(10)
    pdf.close()
    out = f"/tmp/conv_{job_id}.pptx"
    prs.save(out)
    return f"conversions/{user_id}/{job_id}.pptx", out


def _office_to_pdf(src: str, job_id: str, src_fmt: str, user_id: UUID) -> tuple[str, str]:
    out_dir = f"/tmp/conv_{job_id}_out"
    os.makedirs(out_dir, exist_ok=True)

    result = subprocess.run(
        [LIBREOFFICE_BIN, "--headless", "--convert-to", "pdf", "--outdir", out_dir, src],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {result.stderr}")

    converted = next(
        (os.path.join(out_dir, f) for f in os.listdir(out_dir) if f.endswith(".pdf")),
        None,
    )
    if not converted:
        raise HTTPException(status_code=500, detail="Conversion produced no output")

    key = f"conversions/{user_id}/{job_id}.pdf"
    return key, converted


def _mime(fmt: str) -> str:
    return {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "png":  "image/png",
        "jpg":  "image/jpeg",
        "txt":  "text/plain",
        "html": "text/html",
    }.get(fmt, "application/octet-stream")


